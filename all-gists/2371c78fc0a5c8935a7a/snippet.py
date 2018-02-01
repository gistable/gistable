#!/usr/bin/env python2
# coding: utf-8

from pptx import Presentation
import chardet
import sys

def main(filename):

    prs = Presentation(filename)

    # text_runs will be populated with a list of strings,
    # one for each text run in presentation
    text_runs = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                print u''.join(run.text for run in paragraph.runs).encode('utf-8')
                #text = u''.join(run.text.decode(detector.feed(line)['encoding']) for run in paragraph.runs)
                #text_runs.append(text)

if __name__ == "__main__":
    main(sys.argv[1])
