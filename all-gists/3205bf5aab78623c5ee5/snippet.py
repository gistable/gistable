"""
Example of accessing the notes slides of a presentation.
Requires python-pptx 0.5.6 or later.

ryan@ryanday.net
"""

from pptx.util import lazyproperty, Pt
from pptx.parts.slide import BaseSlide, Slide, _SlideShapeTree, _SlidePlaceholders
from pptx.shapes.shape import BaseShape
from pptx.opc.constants import RELATIONSHIP_TYPE as RT, CONTENT_TYPE as CT
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER_TYPE
from pptx.oxml.xmlchemy import BaseOxmlElement, OneAndOnlyOne, ZeroOrOne
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
from pptx.opc.constants import CONTENT_TYPE as CT
from pptx.oxml import register_element_cls
from pptx import content_type_to_part_class_map
from pptx.opc.package import PartFactory


"""
http://msdn.microsoft.com/en-us/library/office/gg278319%28v=office.15%29.aspx
"""


class CT_NotesSlide(BaseOxmlElement):
    """
    ``<p:notes>`` element, root of a notesSlide part
    """
    cSld = OneAndOnlyOne('p:cSld')
    clrMapOvr = ZeroOrOne('p:clrMapOvr', successors=(
        'p:transition', 'p:timing', 'p:extLst'
    ))

    @classmethod
    def new(cls):
        """
        Return a new ``<p:notes>`` element configured as a base slide shape.
        """
        return parse_xml(cls._notes_xml())

    @staticmethod
    def _notes_xml():
        """From http://msdn.microsoft.com/en-us/library/office/gg278319%28v=office.15%29.aspx#sectionSection4
        """
        return (
           '<p:notes %s>\n'
           '   <p:cSld>\n'
           '     <p:spTree>\n'
           '       <p:nvGrpSpPr>\n'
           '         <p:cNvPr id="1"\n'
           '                  name="" />\n'
           '         <p:cNvGrpSpPr />\n'
           '         <p:nvPr />\n'
           '       </p:nvGrpSpPr>\n'
           '       <p:grpSpPr>\n'
           '         <a:xfrm xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" />\n'
           '       </p:grpSpPr>\n'
           '       <p:sp>\n'
           '         <p:nvSpPr>\n'
           '           <p:cNvPr id="2"\n'
           '                    name="" />\n'
           '           <p:cNvSpPr>\n'
           '             <a:spLocks noGrp="1"\n'
           '                        xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" />\n'
           '           </p:cNvSpPr>\n'
           '           <p:nvPr>\n'
           '             <p:ph />\n'
           '           </p:nvPr>\n'
           '         </p:nvSpPr>\n'
           '         <p:spPr />\n'
           '         <p:txBody>\n'
           '           <a:bodyPr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" />\n'
           '           <a:lstStyle xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" />\n'
           '           <a:p xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">\n'
           '             <a:endParaRPr />\n'
           '           </a:p>\n'
           '         </p:txBody>\n'
           '       </p:sp>\n'
           '     </p:spTree>\n'
           '   </p:cSld>\n'
           '   <p:clrMapOvr>\n'
           '     <a:masterClrMapping xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" />\n'
           '  </p:clrMapOvr>\n'
           '</p:notes>\n' % nsdecls('p', 'a', 'r')
        )


class SlideWrapper(Slide):
    """Wrapper class for Slide that provides the .notes() property
    """
    def __init__(self, slide):
        self._slide = slide
        self._notes_slide = None

    def __getattr__(self, attr):
        return self._slide.__getattribute__(attr)

    @lazyproperty
    def notes(self):
        """Return all related notesSlides (from @inl1ne's branch)
        """
        notes_slide = None
        try:
            notes_slide = self.part_related_by(RT.NOTES_SLIDE)
        except KeyError:
            notes_slide = SlideNotes.new(self, self.partname, self.package)
            rId = self.relate_to(notes_slide, RT.NOTES_SLIDE)

        return self.part_related_by(RT.NOTES_SLIDE)

    def create_notes_slide(self):
        """
        """
        self._notes_slide = NotesSlide.new(self, self.partname, self.package)
        rId = self.relate_to(self._notes_slide, RT.NOTES_SLIDE)
        return rId


class NotesSlide(BaseSlide):
    """This class will represent the Part of the notesSlide. Any notes retrieved
    from the presentation slides will be an instance of this class.
    """
    @classmethod
    def new(cls, slide, partname, package):
        notes_slide_elm = CT_NotesSlide.new()
        slide = cls(partname, CT.PML_NOTES_SLIDE, notes_slide_elm, package)
        return slide

    @lazyproperty
    def shapes(self):
        """
        Instance of |_SlideShapeTree| containing sequence of shape objects
        appearing on this slide.
        """
        return _SlideShapeTree(self)

    @lazyproperty
    def placeholders(self):
        """
        Instance of |_SlidePlaceholders| containing sequence of placeholder
        shapes in this slide.
        """
        return _SlidePlaceholders(self)

    def add_multiline_note(self, text):
        for line in text.split('\n'):
            self.add_note(line)

    def add_note(self, text):
        """Add some text to the notesSlide, return paragraph
        that was added or False if no textframes were found
        """
        for shape in self.shapes:
            if shape.has_text_frame and shape.is_placeholder:
                if hasattr(shape, 'ph_type') and shape.ph_type == PP_PLACEHOLDER_TYPE.BODY:
                    return shape.text_frame.add_paragraph()
        return False

    def clear_notes(self):
        """Remove all current notes from the slide
        """
        for shape in self.shapes:
            if shape.has_text_frame:
                shape.text_frame.clear()

    def get_slide_runs(self):
        for shape in self.shapes:
            if shape.has_text_frame and shape.is_placeholder:
                for p in shape.text_frame.paragraphs:
                    for run in p.runs:
                        yield run




# Register the <p:notes> root tag
register_element_cls('p:notes', CT_NotesSlide)

# Add our NotesSlide as a valid Part for the notes Content Type
content_type_to_part_class_map[CT.PML_NOTES_SLIDE] = NotesSlide
PartFactory.part_type_for.update(content_type_to_part_class_map)


"""
slide = self.presentation.slides[1]
notes_slide = SlideWrapper(slide).notes_page()
"""