"""
Extract text from files with common file extensions. (inspried by textract <http://textract.readthedocs.org/en/latest/>)
Modified: Dec 2015
System requirement:
    + catdoc/catppt for doc/ppt extract <http://www.wagner.pp.ru/~vitus/software/catdoc/>
"""
import zipfile
import PyPDF2
from subprocess import Popen, PIPE
from pptx import Presentation
import xlrd
import sys
thismodule = sys.modules[__name__]


SUPPORT_EXTENTIONS = ['docx', 'doc', 'pdf', 'ppt', 'pptx', 'xls', 'xlsx']
try:
    from xml.etree.cElementTree import XML
except ImportError:
    from xml.etree.ElementTree import XML

"""
Module that extract text from MS XML Word document (.docx).
(Inspired by python-docx <https://github.com/mikemaccana/python-docx>)
"""

WORD_NAMESPACE = '{http://schemas.openxmlformats.org/wordprocessingml/2006/main}'
PARA = WORD_NAMESPACE + 'p'
TEXT = WORD_NAMESPACE + 't'


def get_docx_text(path):
    """
    Take the path of a docx file as argument, return the text in unicode.
    """
    document = zipfile.ZipFile(path)
    xml_content = document.read('word/document.xml')
    document.close()
    tree = XML(xml_content)

    paragraphs = []
    for paragraph in tree.getiterator(PARA):
        texts = [node.text
                 for node in paragraph.getiterator(TEXT)
                 if node.text]
        if texts:
            paragraphs.append(''.join(texts))

    return '\n\n'.join(paragraphs)


def get_pdf_text(path):
    pdfFileObj = open(path, 'rb')
    pdfReader = PyPDF2.PdfFileReader(pdfFileObj)
    res = []
    for i in range(pdfReader.numPages):
        pageObj = pdfReader.getPage(i)
        res.append(pageObj.extractText())
    return "\n".join(res)


def get_doc_text(path):
    #cmd = ['antiword', '-m', 'utf-8.txt', path]
    cmd = ['catdoc', '-d', 'utf-8', path]
    try:
        p = Popen(cmd, stdout=PIPE)
        stdout, stderr = p.communicate()
        return stdout.decode('utf-8', 'ignore')
    except:
        return ''


def get_ppt_text(path):
    cmd = ['catppt', '-d', 'utf-8', path]
    try:
        p = Popen(cmd, stdout=PIPE)
        stdout, stderr = p.communicate()
        return stdout.decode('utf-8', 'ignore')
    except:
        return ''


def get_pptx_text(path):
    prs = Presentation(path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs.append(run.text)

    return " ".join(text_runs)


def get_xls_text(path):
    workbook = xlrd.open_workbook(path)
    sheets_name = workbook.sheet_names()
    output = "\n"
    for names in sheets_name:
        worksheet = workbook.sheet_by_name(names)
        num_rows = worksheet.nrows
        num_cells = worksheet.ncols

        for curr_row in range(num_rows):
            #row = worksheet.row(curr_row)
            new_output = []
            for index_col in xrange(num_cells):
                value = worksheet.cell_value(curr_row, index_col)
                if value:
                    if isinstance(value, (int, float)):
                        value = unicode(value)
                    new_output.append(value)
            if new_output:
                output += u' '.join(new_output) + u'\n'
    return output


def get_xlsx_text(path):
    return get_xls_text(path)


def textract(fname, ftype):
    if ftype in SUPPORT_EXTENTIONS:
        return getattr(thismodule, 'get_' + str(ftype) + '_text')(fname)
    else:
        f = open(fname, 'r')
        return f.read()
