##
##  Copy & Paste Tool for images to PowerPoint(.pptx)
##
import pptx
import pptx.util
import glob
import scipy.misc

OUTPUT_TAG = "MY_TAG"

# new
prs = pptx.Presentation()
# open
# prs_exists = pptx.Presentation("some_presentation.pptx")

# default slide width
#prs.slide_width = 9144000
# slide height @ 4:3
#prs.slide_height = 6858000
# slide height @ 16:9
prs.slide_height = 5143500

# title slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
# blank slide
#slide = prs.slides.add_slide(prs.slide_layouts[6])

# set title
title = slide.shapes.title
title.text = OUTPUT_TAG

pic_left  = int(prs.slide_width * 0.15)
pic_top   = int(prs.slide_height * 0.1)
pic_width = int(prs.slide_width * 0.7)

for g in glob.glob("/path/to/your/images*"):
    print g
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    tb = slide.shapes.add_textbox(0, 0, prs.slide_width, pic_top / 2)
    p = tb.textframe.add_paragraph()
    p.text = g
    p.font.size = pptx.util.Pt(14)

    img = scipy.misc.imread(g)
    pic_height = int(pic_width * img.shape[0] / img.shape[1])
    #pic   = slide.shapes.add_picture(g, pic_left, pic_top)
    pic   = slide.shapes.add_picture(g, pic_left, pic_top, pic_width, pic_height)

prs.save("%s.pptx" % OUTPUT_TAG)