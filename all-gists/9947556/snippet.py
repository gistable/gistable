#!/usr/bin/python
# -*- coding: utf-8 -*-

import os
import sys
from cStringIO import StringIO
from argparse import ArgumentParser

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


def is_md(src, f):
    """Returns true if file exists and has a Markdown extension"""
    return os.path.isfile(os.path.join(src, f)) and f.endswith('.md')


def get_pages(src):
    """Returns a list of Markdown filenames in the source directory."""
    return [os.path.join(src, f) for f in os.listdir(src) if is_md(src, f)]


def build_deck(pages):
    """Returns the actual PowerPoint deck."""
    pres = Presentation()
    blank_slidelayout = pres.slide_layouts[6]

    for page in pages:
        slide = pres.slides.add_slide(blank_slidelayout)

        left = top = Inches(0)
        width = Inches(10)
        height = Inches(5.63)

        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width,
                                       height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0, 0, 0)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        textframe = textbox.text_frame

        para = textframe.paragraphs[0]
        run = para.add_run()
        run.font.name = "Consolas"
        run.font.size = Pt(7)
        run.font.color.rgb = RGBColor(255, 255, 255)

        f = open(page)
        run.text = f.read()
        f.close()

    return pres


if __name__ == '__main__':
    parser = ArgumentParser(
        description='My dorky colleague created a presentation using vimdeck.\
                     The very last thing he would want to see is it as a real\
                     PowerPoint deck. So here you go...',
        epilog='This thing returns binary data, so you probably want to pipe\
                the output to a file with a .pptx extension. Oh, and be sure\
                to "pip install python-pptx" first.')
    parser.add_argument('source', type=str, nargs=1, help='Path to source\
        directory containing vimdeck markdown files.')
    args = parser.parse_args()

    source = os.path.realpath(args.source[0])
    deck = build_deck(get_pages(source))

    output = StringIO()
    deck.save(output)
    sys.stdout.write(output.getvalue())
    output.close()
    sys.exit()
